"""ai_text_files.builders_files — assemble final output files.

Notes:
  • create_pdf — renders the document structure as a PDF using reportlab
  • list_outputs lists .pdf files alongside .docx/.pptx/.py

Public API:
    create_docx(topic, model, file_context="", n_charts=None) → str path
    create_pptx(topic, model, file_context="", n_charts=None) → str path
    create_pdf(topic, model, file_context="", n_charts=None)  → str path
    create_python_script(description, model, file_context="") → str path
    list_outputs() — print contents of OUTPUT_DIR
"""

from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGBColor
from pptx.util import Inches as PInches, Pt as PPt

from . import (builders_structure, charts, config, conversation,
               intent_router)


# ─── Color helpers ──────────────────────────────────────────────────
def _hclr(p, hx="1F3864"):
    r, g, b = int(hx[:2], 16), int(hx[2:4], 16), int(hx[4:], 16)
    for run in p.runs:
        run.font.color.rgb = RGBColor(r, g, b)


def _h2r(h):
    h = h.lstrip('#')
    return PRGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))


def _resolve_filename(prompt_text, ext):
    for pat in [
        rf'(?:named?|call(?:\s+it)?|save(?:\s+(?:it|as|to))*|output(?:\s+to)?)\s+'
        rf'(?:as\s+|to\s+|it\s+)?([\w\-]+\.{ext})',
        rf'\b([\w\-]+\.{ext})\b',
    ]:
        m = re.search(pat, prompt_text, re.IGNORECASE)
        if m:
            print(f"  📁 Custom filename: {m.group(1)}")
            return config.OUTPUT_DIR / m.group(1)
    return None


def _fallback_filename(prompt_text, ext):
    safe = re.sub(r'[^\w\s-]', '', prompt_text)[:40].strip().replace(' ', '_')
    return config.OUTPUT_DIR / f"{safe}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"


# ─── F-06: deck framing options ─────────────────────────────────────
def _parse_deck_options(topic, n_content_slides):
    """Return (include_title_slide, include_closer_slide) based on prompt
    text and content-slide count.

    Defaults:
      • Title slide: ON (a deck without a title is unusual)
      • Closer slide: ON only when the deck has more than 5 content slides
        (under 6 slides the Thank You feels like padding)

    User overrides recognized in the topic text:
      "no title slide" / "without title slide" / "no cover slide" → title OFF
      "no closing slide" / "no closer" / "no thank you" / "no thank-you slide"
                                                          → closer OFF
      "with title slide" / "include title slide"           → title ON
      "with closing slide" / "include thank you slide"     → closer ON
    """
    t = (topic or "").lower()

    include_title  = True
    include_closer = n_content_slides > 5

    if re.search(r'\bno\s+(?:title\s+slide|cover\s+slide)\b', t):
        include_title = False
    if re.search(r'\b(?:with|include|add)\s+(?:a\s+)?title\s+slide\b', t):
        include_title = True
    if re.search(r'\bno\s+(?:closing\s+slide|closer\s+slide|'
                 r'thank[\s-]?you(?:\s+slide)?|final\s+slide)\b', t):
        include_closer = False
    if re.search(r'\b(?:with|include|add)\s+(?:a\s+)?(?:closing|'
                 r'thank[\s-]?you)\s+slide\b', t):
        include_closer = True

    return include_title, include_closer


# ─── DOCX ───────────────────────────────────────────────────────────
def create_docx(topic, model, file_context="", n_charts=None):
    if n_charts is None:
        n_charts = intent_router.parse_chart_request(topic)
    print(f"\n📄 Creating DOCX: {topic[:80]}")
    print(f"  Charts requested: {n_charts}")
    chart_plan = (charts.ai_pick_charts(topic, model=model, n_charts=n_charts)
                  if n_charts > 0 else [])
    chart_plan = chart_plan[:n_charts]

    if file_context:
        print("\n  STEP 1 — Reading attached files...")
        analysis = builders_structure.analyze_files_for_task(
            topic, file_context, model=model)
    else:
        analysis = topic

    print("\n  STEP 2 — Structuring document...")
    data = builders_structure.build_docx_structure(
        topic, analysis, model=model, n_chart_sections=len(chart_plan),
        has_files=bool(file_context))

    doc = Document()
    for s in doc.sections:
        s.top_margin = s.bottom_margin = s.left_margin = s.right_margin = Inches(1)

    t = doc.add_heading(data.get("title", topic), level=0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _hclr(t)

    if data.get("subtitle"):
        st = doc.add_paragraph(data["subtitle"])
        st.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if st.runs:
            st.runs[0].italic = True

    doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
    doc.add_paragraph("")

    es = doc.add_heading("Executive Summary", level=1)
    _hclr(es)
    doc.add_paragraph(data.get("executive_summary", ""))

    ci = 0
    for i, sec in enumerate(data.get("sections", [])):
        h = doc.add_heading(sec.get("heading", f"Section {i+1}"), level=1)
        _hclr(h)
        doc.add_paragraph(sec.get("content", ""))
        ctt = sec.get("chart_topic")
        if ctt and ci < len(chart_plan):
            cp = chart_plan[ci]
            use_t = cp["topic"] or ctt
            cpath = charts.generate_chart(
                use_t, chart_type=cp["chart_type"],
                save_path=str(config.OUTPUT_DIR /
                              f"dc_{ci}_{datetime.now().strftime('%f')}.png"))
            doc.add_paragraph("")
            doc.add_picture(cpath, width=Inches(5.5))
            cap = doc.add_paragraph(f"Figure {ci+1}: {use_t}")
            cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if cap.runs:
                cap.runs[0].italic = True
                cap.runs[0].font.size = Pt(9)
            ci += 1

    ch = doc.add_heading("Conclusion", level=1)
    _hclr(ch)
    doc.add_paragraph(data.get("conclusion", ""))

    out = _resolve_filename(topic, "docx") or _fallback_filename(topic, "docx")
    doc.save(str(out))
    print(f"\n✅ DOCX → {out.resolve()}")
    return str(out)


# ─── PDF ────────────────────────────────────────────────────────────
def create_pdf(topic, model, file_context="", n_charts=None):
    """Render the document structure as a PDF using reportlab.
    Mirrors create_docx — same content pipeline, different renderer."""
    try:
        from reportlab.lib.colors import HexColor, black
        from reportlab.lib.pagesizes import LETTER
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import (Image, PageBreak, Paragraph,
                                          SimpleDocTemplate, Spacer)
    except ImportError:
        print("\n❌ reportlab not installed. Re-run the installer cell or "
              "`pip install reportlab`.")
        raise

    if n_charts is None:
        n_charts = intent_router.parse_chart_request(topic)
    print(f"\n📕 Creating PDF: {topic[:80]}")
    print(f"  Charts requested: {n_charts}")
    chart_plan = (charts.ai_pick_charts(topic, model=model, n_charts=n_charts)
                  if n_charts > 0 else [])
    chart_plan = chart_plan[:n_charts]

    if file_context:
        print("\n  STEP 1 — Reading attached files...")
        analysis = builders_structure.analyze_files_for_task(
            topic, file_context, model=model)
    else:
        analysis = topic

    print("\n  STEP 2 — Structuring PDF...")
    data = builders_structure.build_pdf_structure(
        topic, analysis, model=model, n_chart_sections=len(chart_plan),
        has_files=bool(file_context))

    out = _resolve_filename(topic, "pdf") or _fallback_filename(topic, "pdf")

    # Set up document with 1-inch margins
    doc = SimpleDocTemplate(
        str(out), pagesize=LETTER,
        leftMargin=inch, rightMargin=inch,
        topMargin=inch, bottomMargin=inch,
        title=data.get("title", topic)[:120],
        author="ai_text",
    )

    # Styles (navy headings, black body)
    NAVY  = HexColor("#1F3864")
    GREY  = HexColor("#666666")
    base  = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "Title", parent=base["Title"],
        fontName="Helvetica-Bold", fontSize=22,
        textColor=NAVY, alignment=1, spaceAfter=10)
    subtitle_style = ParagraphStyle(
        "Subtitle", parent=base["Normal"],
        fontName="Helvetica-Oblique", fontSize=12,
        textColor=GREY, alignment=1, spaceAfter=6)
    meta_style = ParagraphStyle(
        "Meta", parent=base["Normal"],
        fontName="Helvetica", fontSize=9,
        textColor=GREY, alignment=1, spaceAfter=20)
    h1_style = ParagraphStyle(
        "H1", parent=base["Heading1"],
        fontName="Helvetica-Bold", fontSize=15,
        textColor=NAVY, spaceBefore=14, spaceAfter=8)
    body_style = ParagraphStyle(
        "Body", parent=base["BodyText"],
        fontName="Helvetica", fontSize=10.5, leading=14,
        textColor=black, spaceAfter=8, alignment=4)  # 4 = justify
    caption_style = ParagraphStyle(
        "Caption", parent=base["Normal"],
        fontName="Helvetica-Oblique", fontSize=9,
        textColor=GREY, alignment=1, spaceAfter=12, spaceBefore=4)

    flow = []

    # Title block
    flow.append(Paragraph(data.get("title", topic), title_style))
    if data.get("subtitle"):
        flow.append(Paragraph(data["subtitle"], subtitle_style))
    flow.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')}", meta_style))

    # Executive summary
    flow.append(Paragraph("Executive Summary", h1_style))
    for para in (data.get("executive_summary", "") or "").split("\n\n"):
        if para.strip():
            flow.append(Paragraph(_pdf_escape(para), body_style))

    # Sections (with optional charts)
    ci = 0
    for i, sec in enumerate(data.get("sections", [])):
        flow.append(Paragraph(_pdf_escape(sec.get("heading", f"Section {i+1}")),
                              h1_style))
        for para in (sec.get("content", "") or "").split("\n\n"):
            if para.strip():
                flow.append(Paragraph(_pdf_escape(para), body_style))
        ctt = sec.get("chart_topic")
        if ctt and ci < len(chart_plan):
            cp = chart_plan[ci]
            use_t = cp["topic"] or ctt
            cpath = charts.generate_chart(
                use_t, chart_type=cp["chart_type"],
                save_path=str(config.OUTPUT_DIR /
                              f"pdfc_{ci}_{datetime.now().strftime('%f')}.png"))
            flow.append(Spacer(1, 6))
            try:
                flow.append(Image(cpath, width=5.5*inch, height=3.1*inch))
                flow.append(Paragraph(f"Figure {ci+1}: {_pdf_escape(use_t)}",
                                      caption_style))
            except Exception as e:
                flow.append(Paragraph(
                    f"[chart could not be embedded: {type(e).__name__}: {e}]",
                    caption_style))
            ci += 1

    # Conclusion
    flow.append(Paragraph("Conclusion", h1_style))
    for para in (data.get("conclusion", "") or "").split("\n\n"):
        if para.strip():
            flow.append(Paragraph(_pdf_escape(para), body_style))

    doc.build(flow)
    print(f"\n✅ PDF → {out.resolve()}")
    return str(out)


def _pdf_escape(text):
    """Escape characters that have meaning in reportlab's mini-XML."""
    return (str(text)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;"))


# ─── PPTX ───────────────────────────────────────────────────────────
def create_pptx(topic, model, file_context="", n_charts=None,
                include_title_slide=None, include_closer_slide=None):
    """Build a .pptx from a topic.

    F-06: title and closer ("Thank You") slides are now configurable.
      • include_title_slide / include_closer_slide: explicit True/False
        overrides everything else.
      • If left as None, _parse_deck_options reads the topic for phrases
        like "no closing slide" / "no thank you" and defaults the closer
        OFF for decks with ≤5 content slides.
    """
    if n_charts is None:
        n_charts = intent_router.parse_chart_request(topic)
    print(f"\n📊 Creating PPTX: {topic[:80]}")
    print(f"  Charts requested: {n_charts}")
    chart_plan = (charts.ai_pick_charts(topic, model=model, n_charts=n_charts)
                  if n_charts > 0 else [])
    chart_plan = chart_plan[:n_charts]

    if file_context:
        print("\n  STEP 1 — Reading attached files...")
        analysis = builders_structure.analyze_files_for_task(
            topic, file_context, model=model)
    else:
        analysis = topic

    print("\n  STEP 2 — Structuring slides...")
    data = builders_structure.build_pptx_structure(
        topic, analysis, model=model, n_chart_slides=len(chart_plan),
        has_files=bool(file_context))

    # F-06: decide framing (title slide / closer slide)
    content_slides = data.get("slides", [])
    n_content = len(content_slides)
    auto_title, auto_closer = _parse_deck_options(topic, n_content)
    if include_title_slide  is None: include_title_slide  = auto_title
    if include_closer_slide is None: include_closer_slide = auto_closer
    print(f"  🎬 Framing: title_slide={include_title_slide}  "
          f"closer_slide={include_closer_slide}  "
          f"(content slides: {n_content})")

    DARK, LIGHT = "1F3864", "FFFFFF"
    prs = Presentation()
    prs.slide_width  = PInches(13.33)
    prs.slide_height = PInches(7.5)

    if include_title_slide:
        ts = prs.slides.add_slide(prs.slide_layouts[0])
        ts.background.fill.solid()
        ts.background.fill.fore_color.rgb = _h2r(DARK)
        ts.shapes.title.text = data.get("title", topic)
        for p in ts.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = _h2r(LIGHT)
                r.font.size = PPt(40)
                r.font.bold = True
        if len(ts.placeholders) > 1:
            ts.placeholders[1].text = data.get("subtitle", "")
            for p in ts.placeholders[1].text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = _h2r("CADCFC")
                    r.font.size = PPt(20)

    ci = 0
    for sd in content_slides:
        has_c = bool(sd.get("chart_topic") and ci < len(chart_plan))
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _h2r(LIGHT)
        slide.shapes.title.text = sd.get("title", "")
        for p in slide.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = _h2r(DARK)
                r.font.size = PPt(28)
                r.font.bold = True
        if len(slide.placeholders) > 1:
            bp = slide.placeholders[1]
            if has_c:
                bp.left   = PInches(0.3)
                bp.top    = PInches(1.5)
                bp.width  = PInches(5.5)
                bp.height = PInches(5.5)
            tf = bp.text_frame
            tf.word_wrap = True
            for i, b in enumerate(sd.get("bullets", [])):
                p2 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p2.text = b
                p2.font.size = PPt(13)
        if has_c:
            cp = chart_plan[ci]
            use_t = cp["topic"] or sd["chart_topic"]
            cpath = charts.generate_chart(
                use_t, chart_type=cp["chart_type"],
                save_path=str(config.OUTPUT_DIR /
                              f"pc_{ci}_{datetime.now().strftime('%f')}.png"))
            slide.shapes.add_picture(cpath, PInches(6.2), PInches(1.5), PInches(6.8))
            ci += 1

    if include_closer_slide:
        cs = prs.slides.add_slide(prs.slide_layouts[1])
        cs.background.fill.solid()
        cs.background.fill.fore_color.rgb = _h2r(DARK)
        cs.shapes.title.text = "Thank You"
        for p in cs.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = _h2r(LIGHT)
                r.font.size = PPt(36)
        if len(cs.placeholders) > 1:
            cs.placeholders[1].text = data.get("subtitle", topic)
            for p in cs.placeholders[1].text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = _h2r("CADCFC")
                    r.font.size = PPt(18)

    out = _resolve_filename(topic, "pptx") or _fallback_filename(topic, "pptx")
    prs.save(str(out))
    print(f"\n✅ PPTX → {out.resolve()}")
    return str(out)


# ─── Python script ──────────────────────────────────────────────────
def create_python_script(description, model, file_context=""):
    print(f"\n🐍 Generating: {description[:80]}")
    has_files = bool(file_context)
    if has_files:
        print("\n  STEP 1 — Reading attached files...")
        analysis = builders_structure.analyze_files_for_task(
            description, file_context, model=model)
        print("\n  STEP 2 — Generating code...")
        code = builders_structure.build_python_code(
            description, analysis, model=model, has_files=True)
    else:
        code = builders_structure.build_python_code(
            description, description, model=model, has_files=False)

    out = _resolve_filename(description, "py") or _fallback_filename(description, "py")
    out.write_text(code, encoding='utf-8')
    print(f"\n✅ Script → {out.resolve()}")
    print("─" * 60)
    print(code[:1800] + ("\n[... truncated]" if len(code) > 1800 else ""))
    return str(out)


# ─── List outputs ───────────────────────────────────────────────────
def list_outputs():
    files  = [f for f in sorted(config.OUTPUT_DIR.glob("*.*"))
              if f.suffix not in ('.png',) and f.is_file()]
    charts_pngs = list(config.OUTPUT_DIR.glob("*.png"))
    if files:
        print(f"📁 {config.OUTPUT_DIR.resolve()}\n")
        print(f"  {'Name':<55} {'Size':>8}")
        print("  " + "─" * 65)
        for f in files:
            print(f"  {f.name:<55} {f.stat().st_size/1024:>7.1f} KB")
        if charts_pngs:
            print(f"\n  + {len(charts_pngs)} chart PNG(s)")
    else:
        print("No output files yet.")
    logs = sorted(config.LOGS_DIR.glob("*.txt"))
    if logs:
        print(f"\n📋 Logs in {config.LOGS_DIR}:")
        for lg in logs:
            print(f"  {lg.name:<55} {lg.stat().st_size/1024:>7.1f} KB")
