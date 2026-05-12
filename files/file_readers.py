"""ai_text_files.file_readers — read attached files into prompt-ready strings.

added .pdf handling using pypdf.
"""

from __future__ import annotations

import re
from pathlib import Path

import pandas as pd
from docx import Document
from pptx import Presentation

from . import config


def _read_pdf(path):
    """Extract text from a PDF using pypdf. Returns one string with
    page separators. Falls back to a clear error message if pypdf is
    not installed."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return ("[ERROR reading PDF: pypdf is not installed. "
                "Re-run the installer cell or `pip install pypdf`.]")
    try:
        reader = PdfReader(str(path))
        n_pages = len(reader.pages)
        chunks = []
        for i, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text() or ""
            except Exception as e:
                text = f"[error extracting page {i}: {e}]"
            chunks.append(f"--- Page {i} of {n_pages} ---\n{text.strip()}")
        return "\n\n".join(chunks)
    except Exception as e:
        return f"[ERROR reading PDF {path.name}: {type(e).__name__}: {e}]"


def read_file_by_path(path):
    path = Path(path)
    ext = path.suffix.lower()
    try:
        if ext == ".docx":
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".pptx":
            prs = Presentation(str(path))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                lines.append(para.text)
            return "\n".join(lines)
        elif ext == ".pdf":
            return _read_pdf(path)
        elif ext == ".csv":
            df = pd.read_csv(path)
            return f"Shape: {df.shape}\n\n{df.head(30).to_string()}"
        elif ext == ".xlsx":
            df = pd.read_excel(path)
            return f"Shape: {df.shape}\n\n{df.head(30).to_string()}"
        else:
            return path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"[ERROR reading {path.name}: {e}]"


def build_file_context(file_paths, per_file_limit=15000):
    parts, total = [], 0
    for fp in file_paths:
        p = Path(fp)
        if not p.exists():
            parts.append(f"\n=== FILE: {p.name} ===\n[NOT FOUND: {fp}]")
            continue
        content = read_file_by_path(p)
        if len(content) > per_file_limit:
            content = content[:per_file_limit] + f"\n...[truncated — {len(content):,} total chars]"
        parts.append(f"\n=== FILE: {p.name} ===\n{content}\n=== END: {p.name} ===")
        total += len(content)
        print(f"  📎 {p.name}  ({len(content):,} chars)")
    print(f"  ─ Total: {total:,} chars across {len(file_paths)} file(s)")
    return "\n".join(parts)


def find_local_file(filename):
    candidates = [config.WORK_DIR / filename, config.OUTPUT_DIR / filename]
    try:
        candidates += [sub / filename for sub in config.WORK_DIR.iterdir() if sub.is_dir()]
    except Exception:
        pass
    for p in candidates:
        if p.exists() and p.is_file():
            return p
    return None


def detect_read_request(text):
    m = re.search(r'(?:read|open|load|show|display)\s+([\w.\-]+\.\w+)', text, re.IGNORECASE)
    return (True, m.group(1)) if m else (False, None)
